import pdfplumber
from PIL import Image
import pytesseract
import os
import logging
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown
from pathlib import Path
from ..utils.config_manager import ConfigManager
from ..utils.error_handler import ErrorHandler

class FileProcessor:
    def __init__(self):
        self.config = ConfigManager()
        self._setup_ocr()
    
    def _setup_ocr(self):
        """设置OCR"""
        ocr_path = self.config.get('ocr.tesseract_cmd')
        if ocr_path:
            pytesseract.pytesseract.tesseract_cmd = ocr_path
    
    def process_file(self, file_path):
        """处理文件"""
        try:
            file_path = Path(file_path)
            
            # 检查文件大小
            max_size = self.config.get('file.max_file_size', 20 * 1024 * 1024)  # 默认20MB
            if file_path.stat().st_size > max_size:
                raise ValueError(f"文件大小超过限制: {max_size/1024/1024}MB")
                
            # 根据文件类型处理
            suffix = file_path.suffix.lower()
            
            if suffix == '.txt':
                return self._process_text(file_path)
            elif suffix == '.pdf':
                return self._process_pdf(file_path)
            elif suffix in ['.png', '.jpg', '.jpeg']:
                return self._process_image(file_path)
            elif suffix == '.docx':
                return self._process_docx(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {suffix}")
                
        except Exception as e:
            ErrorHandler.show_error(f"文件处理失败: {str(e)}")
            raise
    
    def _process_text(self, file_path):
        """处理文本文件"""
        content = file_path.read_text(encoding='utf-8')
        
        # 检测模板格式
        paragraphs = content.split('\n\n')  # 两个回车分隔的段落
        if len(paragraphs) >= 2:
            # 第一段作为要求，后面的作为可编辑内容
            requirements = paragraphs[0].strip()
            editable_content = '\n\n'.join(paragraphs[1:]).strip()
            
            return {
                'type': 'template',
                'requirements': requirements,
                'content': editable_content
            }
        
        return {
            'type': 'plain',
            'content': content
        }
    
    def _process_pdf(self, file_path):
        """处理PDF文件"""
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)
    
    def _process_image(self, file_path):
        """处理图片文件"""
        image = Image.open(file_path)
        lang = '+'.join(self.config.get('ocr.languages', ['chi_sim', 'eng']))
        return pytesseract.image_to_string(image, lang=lang)
    
    def _process_docx(self, file_path):
        """处理Word文件"""
        doc = Document(file_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    
    def _process_word(self, file_path):
        """处理Word文档"""
        try:
            doc = Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logging.error(f"Word文档处理失败: {str(e)}")
            raise
            
    def _process_excel(self, file_path):
        """处理Excel文件"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.rows:
                    text.append(' '.join(str(cell.value or '') for cell in row))
            return '\n'.join(text)
        except Exception as e:
            logging.error(f"Excel文件处理失败: {str(e)}")
            raise
            
    def _process_powerpoint(self, file_path):
        """处理PowerPoint文件"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            logging.error(f"PowerPoint文件处理失败: {str(e)}")
            raise
            
    def _process_markdown(self, file_path):
        """处理Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_text = f.read()
            html = markdown.markdown(md_text)
            text = BeautifulSoup(html, 'html.parser').get_text()
            return text
        except Exception as e:
            logging.error(f"Markdown文件处理失败: {str(e)}")
            raise 